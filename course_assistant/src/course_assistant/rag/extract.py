"""Extract plain text from course-material file bytes.

Pure given ``(filename, data)`` — no network, no Drive. Supports PDF, Word
(.docx), PowerPoint (.pptx), and plain-text/code files. Unknown or binary types
return an empty string so the caller can skip them.
"""

from __future__ import annotations

import io

_TEXT_EXTS = frozenset(
    {"txt", "md", "py", "ipynb", "js", "ts", "json", "yaml", "yml", "csv", "sql", "sh"}
)


def _extension(filename: str) -> str:
    _, _, ext = filename.rpartition(".")
    return ext.lower() if ext and ext != filename else ""


def is_extractable(filename: str) -> bool:
    """True if :func:`extract_text` knows how to read this file type."""
    return _extension(filename) in (_TEXT_EXTS | {"pdf", "docx", "pptx"})


def _extract_pdf(data: bytes) -> str:
    from pypdf import PdfReader

    reader = PdfReader(io.BytesIO(data))
    return "\n".join(page.extract_text() or "" for page in reader.pages)


def _extract_docx(data: bytes) -> str:
    from docx import Document as DocxDocument

    document = DocxDocument(io.BytesIO(data))
    return "\n".join(p.text for p in document.paragraphs)


def _extract_pptx(data: bytes) -> str:
    from pptx import Presentation

    prs = Presentation(io.BytesIO(data))
    lines: list[str] = []
    for slide in prs.slides:
        for shape in slide.shapes:
            if shape.has_text_frame:
                lines.extend(p.text for p in shape.text_frame.paragraphs)
    return "\n".join(lines)


def extract_text(filename: str, data: bytes) -> str:
    """Return the plain text of ``data``, dispatching on ``filename``'s extension."""
    ext = _extension(filename)
    if ext == "pdf":
        return _extract_pdf(data)
    if ext == "docx":
        return _extract_docx(data)
    if ext == "pptx":
        return _extract_pptx(data)
    if ext in _TEXT_EXTS:
        return data.decode("utf-8", errors="ignore")
    return ""
